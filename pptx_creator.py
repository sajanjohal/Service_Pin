# Sajan Johal
# V3: Adding command prompt run functionality with argv to read file parameters
# This script will read a CSV file with the names of people who are eligible for a Service Pin
# AND are in attendance of the event
# It will then create a powerpoint presentation with each person's name and award on a slide
# EXPECTED CSV FORMAT:
#   first, last, years, attendance

import sys
import pptx


class Person:
    def __init__(self, first, last, years, section):
        self.first = str(first)
        self.last = str(last)
        self.years = int(years)
        self.section = str(section)


class PresentationBuilder:
    def __init__(self, input_name, output_name):
        self.input_file_name = input_name
        self.output_file_name = output_name
        self.prs = pptx.Presentation('SFO PPT Hero Deck-1 (Arial).pptx')
        self.layout = self.prs.slide_masters[1].slide_layouts[16]
        self.person_list = []

    @property
    def xml_slides(self):
        return self.prs.slides._sldIdLst

    # This function takes in a csv file already opened and returns a list of Person objects.
    # The CSV file is assumed to already be narrowed down to people who are eligible and in attendance
    # file handler -> list of Person objects

    def create_person_list(self):
        try:
            f_in = open(self.input_file_name)
        except FileNotFoundError:
            print('ERROR: Could not find given CSV file.')
            sys.exit(1)
        f_in.readline()
        for aline in f_in:
            line_list = aline.strip().split(',')
            self.person_list.append(Person(line_list[0], line_list[1], line_list[2], line_list[4]))
        f_in.close()
        self.person_list.sort(key=lambda r: r.section)

    # Takes in a list of Person objects and creates a powerpoint with each Person object occupying a new slide
    # List, str -> None

    def create_ppt(self):
        self.clear_slides()
        for person in self.person_list:
            self.add_slide(person)

        try:
            self.prs.save(self.output_file_name)
        except PermissionError:
            print('ERROR: Please close file before attempting to save to it.')
            sys.exit(1)

    def add_slide(self, person):
        slide = self.prs.slides.add_slide(self.layout)
        # Adding name to slide in size 72 font
        name_text_frame = slide.shapes[1].text_frame
        name_text_frame.clear()
        p_name = name_text_frame.paragraphs[0]
        run_name = p_name.add_run()
        run_name.text = person.first + ' ' + person.last
        run_name.font.size = pptx.util.Pt(72)
        # Adding years to slide in size 54 font
        year_text_frame = slide.shapes[0].text_frame
        year_text_frame.clear()
        p_year = year_text_frame.paragraphs[0]
        run_year = p_year.add_run()
        run_year.text = str(person.years) + ' Years'
        run_year.font.size = pptx.util.Pt(54)
        # Adding Section to slide in size 54 font
        section_text_frame = slide.shapes[2].text_frame
        section_text_frame.clear()
        p_sect = section_text_frame.paragraphs[0]
        run_sect = p_sect.add_run()
        run_sect.text = str(person.section)
        run_sect.font.size = pptx.util.Pt(54)

    # This function takes a prs object, goes through every slide in it and deletes it
    # Presentation -> None

    def clear_slides(self):
        for slide in self.prs.slides:
            self.delete_slide(slide)

    # This function takes a Presentation object and deletes the individual slide passed
    # Presentation, slide -> None

    def delete_slide(self, slide):
        # Make a dictionary with the necessary information
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(self.prs.slides._sldIdLst)}
        slide_id = slide.slide_id
        self.prs.part.drop_rel(id_dict[slide_id][1])
        del self.prs.slides._sldIdLst[id_dict[slide_id][0]]


def main():
    file_arg = sys.argv[1:]
    input_csv = file_arg[0]
    output_pptx = file_arg[1]
    prs = PresentationBuilder(input_csv, output_pptx)
    prs.create_person_list()
    prs.create_ppt()


if __name__ == '__main__':
    main()
