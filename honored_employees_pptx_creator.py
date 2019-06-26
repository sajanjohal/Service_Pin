# Sajan Johal
# V1: WORKS, BUT pptx merging is buggy
# This script will read a CSV file with the names of people who are eligible for a Service Pin
# AND are in attendance of the event
# It will then create a powerpoint presentation with each person's name and award on a slide
# EXPECTED CSV FORMAT:
#   first, last, years, attendance

import sys

import pptx
from pptx.util import Pt


class Person:
    def __init__(self, first, last, years):
        self.first = str(first)
        self.last = str(last)
        self.years = int(years)


def create_ppt(person_list, out_file):
    # If change in output pptx layout is wanted, change the example pptx the new pptx is created from in the line
    # one below. If you want to use a different layout w/in the same example pptx, change the layout path found two
    # lines below.
    prs = pptx.Presentation('SFO PPT Hero Deck-4 (Arial).pptx')
    commemorate_layout = prs.slide_masters[0].slide_layouts[0]
    clear_slides(prs)
    for person in person_list:
        add_slide(person, prs, commemorate_layout)
    try:
        prs.save(out_file)
    except PermissionError:
        print('ERROR: Please close the .pptx file you are attempting to save to before trying again.')
        sys.exit(1)


# This function takes in a csv file handler, reads it, and returns a list of Person objects of everyone who will
# get honored at the ceremony
# File handler -> List of Person objects

def create_person_list(f_in):
    person_list = []
    f_in.readline()
    for aline in f_in:
        line_list = aline.strip().split(',')
        person_list.append(Person(line_list[0], line_list[1], line_list[2]))
    return person_list


# This function takes a Person and Presentation object in and adds a slide to prs with the Person information
# Person, Presentation, layout -> None

def add_slide(person, prs, layout):
    slide = prs.slides.add_slide(layout)
    # Adding name to slide in size 72 font
    name_text_frame = slide.shapes[0].text_frame
    name_text_frame.clear()
    p_name = name_text_frame.paragraphs[0]
    run_name = p_name.add_run()
    run_name.text = person.first.upper() + ' ' + person.last.upper()
    run_name.font.size = Pt(54)
    # Adding years to slide in size 54 font
    year_text_frame = slide.shapes[1].text_frame
    year_text_frame.clear()
    p_year = year_text_frame.paragraphs[0]
    run_year = p_year.add_run()
    run_year.text = str(person.years) + ' YEARS AT SFO'
    run_year.font.size = Pt(36)


def main(in_file_name, out_file_name):
    try:
        f_in = open(in_file_name)
    except FileNotFoundError:
        print('ERROR: Could not find file.')
        sys.exit(1)
    person_list = create_person_list(f_in)
    create_ppt(person_list, out_file_name)


# This function takes a prs object, goes through every slide in it and deletes it
# Presentation -> None

def clear_slides(prs):
    for slide in prs.slides:
        delete_slide(prs, slide)


# This function takes a Presentation object and deletes the individual slide passed
# Presentation, slide -> None

def delete_slide(prs, slide):
    # Make a dictionary with the necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


if __name__ == '__main__':
    INPUT_CSV = 'ppt_test_1.csv'
    OUTPUT_PPTX = 'ppt_test_1_v1.pptx'
    main(INPUT_CSV, OUTPUT_PPTX)
